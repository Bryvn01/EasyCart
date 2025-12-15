from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Design Helper Functions ---
    def set_background(slide):
        # A clean, neutral background (very light gray)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = title_text
        subtitle.text = subtitle_text

        # Styling
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
        title.text_frame.paragraphs[0].font.bold = True

        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

    def add_content_slide(prs, title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True

        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20) # Slightly smaller to fit more detail
            p.space_after = Pt(12)
            p.level = 0

    def add_section_header(prs, text):
        slide_layout = prs.slide_layouts[2] # Section Header
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # --- SLIDE 1: Title ---
    add_title_slide(prs, "EasyCart: Enterprise-Grade E-Commerce Platform", "Final Year Project Submission\n\nName: [Your Name]\nSupervisor: [Supervisor Name]\nYear: 2025\nRepository: github.com/bryvn01/easycart")

    # --- SLIDE 2: Introduction ---
    add_content_slide(prs, "Project Overview", [
        "EasyCart is a full-stack e-commerce solution architected for the Kenyan market.",
        "It addresses the gap in affordable, locally-tailored digital commerce tools.",
        "Core Value Proposition:",
        "1. Localized Payments: Seamless M-Pesa integration (STK Push).",
        "2. Performance: Mobile-first React frontend with optimistic UI updates.",
        "3. Scalability: Microservices-ready architecture using Django & PostgreSQL.",
        "4. Security: Enterprise-standard JWT Auth & Role-Based Access Control (RBAC)."
    ])

    # --- SLIDE 3: Problem Statement ---
    add_content_slide(prs, "Problem Statement & Context", [
        "What challenges does EasyCart solve?",
        "1. Payment Fragmentation: Global platforms (Shopify/WooCommerce) struggle with native Mobile Money integration.",
        "2. Technical Complexity: Small businesses need 'zero-config' tools, not complex CMS setups.",
        "3. Performance on Low-Bandwidth: Many users are on mobile data; heavy platforms load slowly.",
        "4. Trust & Security: Managing user data and payments requires strict compliance standards."
    ])

    # --- SLIDE 4: System Architecture (High Level) ---
    add_content_slide(prs, "System Architecture", [
        "[ACTION: Insert diagram from 'ARCHITECTURE_DIAGRAM.md' here]",
        "",
        "Decoupled Monolithic Architecture:",
        "• Frontend: React 18 SPA (Single Page Application) consuming REST APIs.",
        "• Backend: Django REST Framework serving JSON data.",
        "• Database: PostgreSQL 14+ for relational integrity.",
        "• Caching: Redis 7.0 for session management and API response caching.",
        "• Media: Cloudinary CDN for optimized image delivery."
    ])

    # --- SLIDE 5: Technical Stack (Detailed) ---
    add_content_slide(prs, "Technology Stack", [
        "Frontend Engineering:",
        "• React 18 + Hooks for state management.",
        "• TailwindCSS for utility-first, responsive styling.",
        "• Axios for intercepted HTTP requests (handles token refresh automatically).",
        "• React Query for server-state synchronization.",
        "",
        "Backend Engineering:",
        "• Django 5.2.7 + DRF 3.16.1.",
        "• SimpleJWT for secure, stateless authentication.",
        "• Celery 5.5.3 for asynchronous tasks (emails, order processing).",
        "• Gunicorn + Whitenoise for production serving."
    ])

    # --- SLIDE 6: Authentication & Security ---
    add_content_slide(prs, "Security & Authentication", [
        "Implementing defense-in-depth principles:",
        "1. JWT (JSON Web Tokens):",
        "   - Access Token (Short-lived, 5 mins).",
        "   - Refresh Token (Long-lived, 24 hours, HttpOnly cookie).",
        "2. Role-Based Access Control (RBAC):",
        "   - Superadmin, Manager, Editor, Viewer roles.",
        "3. Two-Factor Authentication (2FA): TOTP-based for admin accounts.",
        "4. Secure Headers: HSTS, XSS Protection, Content Security Policy implemented."
    ])

    # --- SLIDE 7: Payment Architecture (M-Pesa) ---
    add_content_slide(prs, "Payment Gateway Integration", [
        "[ACTION: Insert 'Payment Flow Diagram' from 'PAYMENT_ARCHITECTURE.md']",
        "",
        "M-Pesa Daraja API Implementation:",
        "• STK Push (Lipa na M-Pesa Online): Real-time prompt on user's phone.",
        "• Asynchronous Callbacks: Backend listens for M-Pesa confirmation via webhook.",
        "• Transaction Verification: Automatic reconciliation of payment status.",
        "• Fallbacks: Support for Stripe and PayPal for international customers."
    ])

    # --- SLIDE 8: Admin Dashboard Features ---
    add_content_slide(prs, "Admin Dashboard & Analytics", [
        "[ACTION: Insert Screenshot of Admin Dashboard (Charts/Graphs)]",
        "",
        "A dedicated React application for business logic:",
        "• Real-time Visualization: Sales trends, revenue vs. targets.",
        "• Inventory Management: Low-stock alerts and bulk updates.",
        "• Order Workflow: Track status from 'Pending' -> 'Processing' -> 'Delivered'.",
        "• Product Management: Image uploads via Cloudinary, rich text descriptions."
    ])

    # --- SLIDE 9: Database Schema ---
    add_content_slide(prs, "Data Modeling (PostgreSQL)", [
        "Normalized Relational Schema:",
        "• Users: Extended AbstractUser with phone & role fields.",
        "• Products: Indexed by category, slug, and price for fast formulation.",
        "• Orders: Linked to User and Payment tables.",
        "• OrderItems: Snapshot of price at time of purchase (Critical for audit trails).",
        "",
        "[ACTION: Insert 'Database Schema' diagram from 'ARCHITECTURE_VISUAL.md']"
    ])

    # --- SLIDE 10: Implementation Challenges ---
    add_content_slide(prs, "Key Implementation Challenges", [
        "1. CORS & Security:",
        "   - Issue: Communication blocked between localhost:3000 and localhost:8000.",
        "   - Solution: Configured `django-cors-headers` with whitelist for dev/prod origins.",
        "2. Image Optimization:",
        "   - Issue: Large images slowing down mobile loading.",
        "   - Solution: Implemented Cloudinary auto-format (f_auto) and quality (q_auto).",
        "3. State Management:",
        "   - Issue: Cart state lost on refresh.",
        "   - Solution: Persisted Cart context to localStorage + synchronized with Backend API."
    ])

    # --- SLIDE 11: Conclusion ---
    add_content_slide(prs, "Conclusion & Future Roadmap", [
        "Conclusion:",
        "EasyCart demonstrates a production-ready architecture that balances user experience with robust backend security.",
        "",
        "Future Roadmap:",
        "• Microservices: Split 'Auth', 'Product', and 'Order' services for scaling.",
        "• PWA (Progressive Web App): Offline functionality for patchy network areas.",
        "• AI Recommendations: Collaborative filtering based on user purchase history."
    ])

    # --- SLIDE 12: Q&A ---
    add_content_slide(prs, "Thank You", [
        "Questions & Discussion",
        "",
        "Live Demo: https://easycart-frontend-wj9x.onrender.com/",
        "Docs & Code: github.com/bryvn01/easycart"
    ])

    output_file = "EasyCart_Final_Project_Presentation_v2.pptx"
    prs.save(output_file)
    print(f"Successfully generated: {output_file}")

if __name__ == "__main__":
    create_presentation()
